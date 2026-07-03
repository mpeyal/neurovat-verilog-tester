"""Build NeuroVAT_Demo.pptx from the live screenshots in results/deck/.

Pure python-pptx; no external skill required.  Run from repo root:
    python tools/build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(ROOT, "results", "deck")
OUT = os.path.join(ROOT, "NeuroVAT_Demo.pptx")

# ---- palette (matches the app's dark UI) -------------------------------
BG = RGBColor(0x0E, 0x17, 0x26)        # deep navy
BG2 = RGBColor(0x14, 0x20, 0x33)       # panel
INK = RGBColor(0xF4, 0xF7, 0xFB)       # near-white
SUB = RGBColor(0x9A, 0xAA, 0xC0)       # muted gray-blue
AMBER = RGBColor(0xF5, 0xA6, 0x23)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
GREEN = RGBColor(0x49, 0xD0, 0x7A)

EMU = 914400
SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]
IMG_RATIO = 975 / 1664.0


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                           prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=4, font="Segoe UI"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space); p.space_before = Pt(0)
        if isinstance(item, tuple):
            txt, opt = item
        else:
            txt, opt = item, {}
        # bullet support
        bullet = opt.get("bullet", False)
        if bullet:
            _set_bullet(p)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opt.get("size", size))
        r.font.bold = opt.get("bold", bold)
        r.font.color.rgb = opt.get("color", color)
        r.font.name = opt.get("font", font)
    return tb


def _set_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set('indent', '-228600'); pPr.set('marL', '228600')
    bu = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
    buf = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    for tag in ('a:buChar', 'a:buNone', 'a:buAutoNum'):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    pPr.append(buf); pPr.append(bu)


def accent(s, x, y, w=2.2, color=AMBER, h=0.06):
    box(s, x, y, w, h, fill=color)


def picture(s, path, x, y, w):
    h = w * IMG_RATIO
    # subtle frame
    box(s, x - 0.04, y - 0.04, w + 0.08, h + 0.08, fill=BG2)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return h


def header(s, kicker, title):
    text(s, 0.7, 0.42, 11.9, 0.4, [(kicker, {"size": 13, "color": TEAL,
         "bold": True})])
    text(s, 0.7, 0.74, 11.9, 0.8, [(title, {"size": 30, "color": INK,
         "bold": True})])
    accent(s, 0.72, 1.46, 2.0, AMBER)


def content_slide(kicker, title, img, bullets, foot=None):
    s = slide()
    header(s, kicker, title)
    iw = 8.55
    picture(s, os.path.join(DECK, img), 0.7, 1.78, iw)
    # bullets panel on the right
    bx = 9.55
    box(s, bx - 0.1, 1.78, SW - bx - 0.6 + 0.1, 4.9, fill=BG2)
    runs = [(b, {"bullet": True, "size": 15, "color": INK, "space": 10})
            for b in bullets]
    text(s, bx + 0.18, 2.05, SW - bx - 0.95, 4.4, runs)
    if foot:
        text(s, 3.1, 7.0, 8.5, 0.35, [(foot, {"size": 11, "color": SUB})],
             align=PP_ALIGN.CENTER)
    page(s)
    return s


_pageno = [0]


def page(s):
    _pageno[0] += 1
    text(s, 12.2, 7.0, 0.9, 0.35, [(f"{_pageno[0]:02d}", {"size": 11,
         "color": SUB})], align=PP_ALIGN.RIGHT)
    text(s, 0.7, 7.0, 5, 0.35, [("NeuroVAT", {"size": 11, "color": SUB,
         "bold": True})])


# ======================================================================
# 1. TITLE
# ======================================================================
s = slide()
box(s, 0, 0, SW, 0.12, fill=AMBER)
text(s, 0.9, 2.35, 11.5, 1.3, [("NeuroVAT", {"size": 66, "color": INK,
     "bold": True})])
text(s, 0.95, 3.55, 11.5, 0.7, [("Neuromorphic Synapse Modeling — without Cadence",
     {"size": 26, "color": AMBER, "bold": True})])
text(s, 0.95, 4.35, 11.0, 1.4, [
    ("Behavioral Python twins of ECFET / FeFET Verilog-A devices, a live stimulus "
     "designer, a spiking-network trainer on a memristive crossbar, and an embedded "
     "Claude agent — all in one desktop app.", {"size": 16, "color": SUB})])
accent(s, 0.97, 4.18, 3.2, TEAL)
text(s, 0.95, 6.5, 11, 0.4, [("Live demo deck · screenshots captured from the running app",
     {"size": 12, "color": SUB})])
box(s, 0, 7.38, SW, 0.12, fill=BG2)

# ======================================================================
# 2. WHAT IS IT (text)
# ======================================================================
s = slide()
header(s, "OVERVIEW", "What is NeuroVAT?")
# two panels
box(s, 0.7, 1.95, 5.75, 4.6, fill=BG2)
text(s, 1.0, 2.2, 5.2, 0.5, [("The problem", {"size": 19, "color": AMBER,
     "bold": True})])
text(s, 1.0, 2.85, 5.25, 3.6, [
    ("Iterating on neuromorphic synapse devices in Cadence/Spectre is slow: a license, "
     "a compile, a netlist and a long transient for every parameter tweak.",
     {"size": 15, "color": INK, "space": 12}),
    ("Exploring learning behaviour (LTP/LTD, STDP, network accuracy) on top of that loop "
     "is painful and hard to visualise.", {"size": 15, "color": INK})])
box(s, 6.85, 1.95, 5.75, 4.6, fill=BG2)
text(s, 7.15, 2.2, 5.2, 0.5, [("The approach", {"size": 19, "color": TEAL,
     "bold": True})])
text(s, 7.15, 2.85, 5.25, 3.6, [
    ("Fast Python “twins” of each .va device — same physics, no EDA loop.",
     {"size": 15, "color": INK, "bullet": True, "space": 11}),
    ("Design any stimulus and see I / R / G respond live.",
     {"size": 15, "color": INK, "bullet": True, "space": 11}),
    ("Train a spiking LIF network on a crossbar of those devices and watch it learn.",
     {"size": 15, "color": INK, "bullet": True, "space": 11}),
    ("Edit the Verilog-A in-app; ask the built-in Claude agent to explain or tune it.",
     {"size": 15, "color": INK, "bullet": True})])
page(s)

# ======================================================================
# 3..N content slides
# ======================================================================
content_slide(
    "STIMULUS", "Signal Designer", "01_signal_designer.png",
    ["10+ generators: single spike, pulse train, LTP/LTD, paired-pulse, Poisson, "
     "burst, STDP pairs, staircase, custom CSV.",
     "Current- or voltage-drive with pA → A / mV → V unit scaling.",
     "Live waveform preview before you run.",
     "One click sweeps every selected device twin."],
    foot="Default LTP/LTD train shown — 300 potentiating then 300 depressing pulses.")

content_slide(
    "DEVICE RESPONSE", "Results — linked I / R / G", "02_results.png",
    ["Three time-locked plots: drive current, memory resistance R_mem, and synaptic "
     "conductance G.",
     "Shared time axis with interactive zoom, pan and A/B probe markers.",
     "Potentiation vs. depression branches are clearly separated.",
     "Export CSV for a direct Spectre overlay."])

content_slide(
    "ANALYSIS", "Per-pulse LTP / LTD curve", "03_analysis.png",
    ["Retained conductance (or resistance) after each pulse vs. pulse number.",
     "Reproduces the paper's Fig. 3c potentiation/depression staircase.",
     "≥ 250 distinguishable states per ramp, near-linear (R² ≈ 1.0).",
     "Symmetric LTP and LTD branches."])

content_slide(
    "PLASTICITY", "STDP learning window", "04_stdp.png",
    ["Sweeps pre/post spike timing Δt and measures ΔG.",
     "Anti-symmetric window: causal (pre→post) potentiates, anti-causal depresses.",
     "Decay constants come straight from device geometry (τ₁ 22 ms, τ₂ 315 ms, τ₃ 19 s).",
     "Logarithmic Δt sampling — dense on the fast rise, out to the slow tail."])

content_slide(
    "NETWORK", "Neuro Trainer — live crossbar", "05_trainer_weights.png",
    ["A spiking LIF network whose synapses ARE the device twins, on a real crossbar.",
     "Live conductance heatmap (rows = output neurons, cols = inputs).",
     "Device-local STDP or surrogate-gradient (BPTT) learning rules.",
     "Weights update in real time as the network trains."])

# two-image learning-dynamics slide
s = slide()
header(s, "NETWORK", "Neuro Trainer — learning dynamics")
picture(s, os.path.join(DECK, "06_trainer_allweights.png"), 0.7, 1.85, 6.0)
picture(s, os.path.join(DECK, "09_trainer_activity.png"), 7.0, 1.85, 5.6)
text(s, 0.7, 5.7, 6.0, 0.4, [("Every synapse's weight trajectory over training",
     {"size": 13, "color": SUB})])
text(s, 7.0, 5.7, 5.6, 0.4, [("Spike raster + live train / test accuracy curve",
     {"size": 13, "color": SUB})])
text(s, 0.7, 6.25, 11.9, 0.7, [
    ("Watch competition and convergence directly: weight traces fan out and settle "
     "while accuracy climbs to 100% — every signal is live, not a mock-up.",
     {"size": 14, "color": INK})])
page(s)

content_slide(
    "RESULTS", "100% accuracy on XOR", "08_trainer_metrics.png",
    ["Trained ECFET-v3 crossbar solving the non-linear XOR task.",
     "Confusion matrix is fully diagonal — every pattern classified correctly.",
     "Per-class precision / recall / F1 all = 1.000.",
     "Held-out test set + live train/test accuracy curve during training."],
    foot="Loaded from saved_model/ecfetv3_xor_neuro_model.json, then re-trained & evaluated live.")

content_slide(
    "TOOLING", "Verilog-A source & embedded agent", "10_verilog_source.png",
    ["Workspace auto-scans every .va file (ECFET v1/v2/v3, FeFET …).",
     "Editable device parameters with live re-simulation.",
     "In-app Verilog-A editor with one-click write-back.",
     "Claude agent (right) explains devices, writes spike patterns, tunes the network."])

# ======================================================================
# FINAL: validation + getting started
# ======================================================================
s = slide()
header(s, "VALIDATION & START", "Paper-validated — and one command to run")
box(s, 0.7, 1.95, 6.5, 4.7, fill=BG2)
text(s, 1.0, 2.18, 6.0, 0.5, [("Matches Sharbati et al., Adv. Mater. 2018",
     {"size": 17, "color": TEAL, "bold": True})])
rows = [
    ("Fig.3a retained ΔR′", "−10 Ω", "−9.8 Ω"),
    ("Fig.3c distinct states", "> 250", "≥ 250 / ramp"),
    ("Fig.3c linearity", "R² = 0.994", "R² ≈ 1.000"),
    ("Fig.4b τ₁ / τ₂ / τ₃", "22 / 315 ms / 19 s", "22.2 / 312.5 / 19.0"),
    ("Retention (13 h)", "3.2 %", "2.83 %"),
]
y = 2.78
text(s, 1.0, y, 6.0, 0.32, [("observable          paper          model",
     {"size": 12, "color": SUB, "bold": True})])
y += 0.42
for a, b, c in rows:
    text(s, 1.0, y, 3.0, 0.32, [(a, {"size": 13, "color": INK})])
    text(s, 4.05, y, 1.7, 0.32, [(b, {"size": 13, "color": SUB})])
    text(s, 5.55, y, 1.6, 0.32, [(c, {"size": 13, "color": GREEN, "bold": True})])
    y += 0.42
text(s, 1.0, y + 0.05, 6.0, 0.4, [("Every electrical synaptic observable reproduced "
     "to figure-reading accuracy.", {"size": 12, "color": SUB})])

box(s, 7.55, 1.95, 5.05, 4.7, fill=BG2)
text(s, 7.85, 2.18, 4.5, 0.5, [("Getting started", {"size": 17, "color": AMBER,
     "bold": True})])
text(s, 7.85, 2.85, 4.5, 0.4, [("pip install -r requirements.txt",
     {"size": 14, "color": INK, "font": "Consolas"})])
text(s, 7.85, 3.35, 4.5, 0.4, [("python run_gui.py",
     {"size": 16, "color": GREEN, "font": "Consolas", "bold": True})])
text(s, 7.85, 4.1, 4.5, 2.3, [
    ("Auto-loads every .va in the folder.", {"size": 14, "color": INK, "bullet": True,
     "space": 9}),
    ("Click Neuro Trainer → Load a saved model → Train.", {"size": 14, "color": INK,
     "bullet": True, "space": 9}),
    ("Export CSV for Spectre cross-checks anytime.", {"size": 14, "color": INK,
     "bullet": True})])
page(s)

prs.save(OUT)
print("saved", OUT, "—", len(prs.slides._sldIdLst), "slides")
